"""PPTX parser — slide-by-slide text, including speaker notes and table content."""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from ai_fde.core.models import ParsedDocument, ParsedSection


class PptxParser:
    def supports(self, filename: str) -> bool:
        return filename.lower().endswith(".pptx")

    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            presentation = Presentation(io.BytesIO(content))
        except PackageNotFoundError as exc:
            raise ValueError(f"{filename}: not a valid .pptx file") from exc

        sections: list[ParsedSection] = []
        for index, slide in enumerate(presentation.slides, start=1):
            heading = _slide_title(slide) or f"Slide {index}"
            body_parts = [p for p in _slide_shape_text(slide) if p]
            notes = _slide_notes(slide)
            if notes:
                body_parts.append(f"Speaker notes: {notes}")
            sections.append(ParsedSection(heading=heading, text="\n\n".join(body_parts)))

        full_text = "\n\n".join(f"## {s.heading}\n\n{s.text}" for s in sections)

        return ParsedDocument(
            filename=filename,
            type="pptx",
            page_or_slide_count=len(presentation.slides),
            sections=sections,
            char_count=len(full_text),
            content=full_text,
        )


def _slide_title(slide) -> str | None:
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        text = slide.shapes.title.text_frame.text.strip()
        return text or None
    return None


def _slide_shape_text(slide) -> list[str]:
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows = [
                " | ".join(cell.text.strip() for cell in row.cells) for row in table.rows
            ]
            parts.append("\n".join(rows))
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)
    return parts


def _slide_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    notes_frame = slide.notes_slide.notes_text_frame
    return notes_frame.text.strip() if notes_frame is not None else ""
